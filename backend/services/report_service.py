"""Engineering report generation service (Phase 10)."""

import io
import json
from typing import Any

from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
from pptx import Presentation
from pptx.util import Pt
from sqlalchemy.orm import Session

from models.investigation import InvestigationRun
from services.errors import InvestigationNotFoundError


class ReportService:
    def __init__(self, db: Session) -> None:
        self._db = db

    def generate_report(self, investigation_id: int, output_format: str = "pdf") -> bytes:
        investigation = self._db.query(InvestigationRun).filter(InvestigationRun.id == investigation_id).one_or_none()
        if investigation is None:
            raise InvestigationNotFoundError(f"Investigation run {investigation_id} was not found")

        state = self._load_state(investigation)
        if output_format == "pptx":
            return self._build_presentation(investigation, state)
        return self._build_pdf(investigation, state)

    def _load_state(self, investigation: InvestigationRun) -> dict[str, Any]:
        if not investigation.state_json:
            return {}

        try:
            parsed = json.loads(investigation.state_json)
            if isinstance(parsed, dict):
                return parsed
        except json.JSONDecodeError:
            pass
        return {}

    def _build_pdf(self, investigation: InvestigationRun, state: dict[str, Any]) -> bytes:
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=LETTER, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72)

        styles = {
            "title": ParagraphStyle("Title", fontSize=20, leading=24, spaceAfter=14, alignment=1),
            "heading": ParagraphStyle("Heading", fontSize=14, leading=18, spaceAfter=10, spaceBefore=12),
            "normal": ParagraphStyle("Normal", fontSize=11, leading=15),
            "bullet": ParagraphStyle("Bullet", parent=ParagraphStyle("Normal"), bulletIndent=12, leftIndent=18, fontSize=11, leading=15),
        }

        story = [Paragraph("Engineering Investigation Report", styles["title"])]
        story.append(Paragraph(f"Investigation ID: {investigation.id}", styles["normal"]))
        story.append(Paragraph(f"Upload ID: {investigation.upload_id}", styles["normal"]))
        story.append(Paragraph(f"Status: {investigation.status}", styles["normal"]))
        story.append(Spacer(1, 0.2 * inch))

        summary = state.get("summary") or "Summary not available."
        story.append(Paragraph("Executive Summary", styles["heading"]))
        story.append(Paragraph(summary, styles["normal"]))
        story.append(Spacer(1, 0.2 * inch))

        story.append(Paragraph("Key Findings", styles["heading"]))
        story.extend(self._build_pdf_findings(state, styles))
        story.append(Spacer(1, 0.2 * inch))

        story.append(Paragraph("Recommendations", styles["heading"]))
        recommendations = state.get("recommendations") or []
        if recommendations:
            for recommendation in recommendations:
                story.append(Paragraph(str(recommendation), styles["bullet"]))
        else:
            story.append(Paragraph("No formal recommendations were generated.", styles["normal"]))

        doc.build(story)
        buffer.seek(0)
        return buffer.read()

    def _build_pdf_findings(self, state: dict[str, Any], styles: dict[str, ParagraphStyle]) -> list[Paragraph]:
        findings: list[Paragraph] = []
        anomalies = state.get("anomalies") or []
        if anomalies:
            findings.append(Paragraph("Detected anomaly patterns:", styles["normal"]))
            for anomaly in anomalies:
                if isinstance(anomaly, dict):
                    findings.append(Paragraph(", ".join(f"{k.replace('_', ' ')}" for k, v in anomaly.items() if v), styles["bullet"]))
        else:
            findings.append(Paragraph("No anomaly patterns detected.", styles["normal"]))

        incidents = state.get("incidents") or []
        if incidents:
            findings.append(Paragraph("Relevant historical incidents:", styles["normal"]))
            for incident in incidents:
                if isinstance(incident, dict):
                    findings.append(
                        Paragraph(
                            f"Incident {incident.get('incident_id')}: similarity={incident.get('similarity')}"  # type: ignore[arg-type]
                            , styles["bullet"]
                        )
                    )
        else:
            findings.append(Paragraph("No relevant historical incidents were matched.", styles["normal"]))

        root_causes = state.get("root_causes") or []
        if root_causes:
            findings.append(Paragraph("Root cause hypotheses:", styles["normal"]))
            for cause in root_causes:
                if isinstance(cause, dict):
                    findings.append(
                        Paragraph(
                            f"{cause.get('cause')} ({cause.get('confidence', 0)}% confidence)", styles["bullet"]
                        )
                    )
        else:
            findings.append(Paragraph("No root cause hypotheses were identified.", styles["normal"]))

        return findings

    def _build_presentation(self, investigation: InvestigationRun, state: dict[str, Any]) -> bytes:
        presentation = Presentation()
        title_slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = "Engineering Investigation Report"
        slide.placeholders[1].text = f"Investigation {investigation.id} • Upload {investigation.upload_id}"

        self._add_bullet_slide(
            presentation,
            "Executive Summary",
            [state.get("summary") or "Summary not available."],
        )

        findings = []
        anomalies = state.get("anomalies") or []
        if anomalies:
            findings.append("Detected anomaly patterns:")
            for anomaly in anomalies:
                if isinstance(anomaly, dict):
                    findings.append(", ".join(f"{k.replace('_', ' ')}" for k, v in anomaly.items() if v))
        else:
            findings.append("No anomaly patterns detected.")

        incidents = state.get("incidents") or []
        if incidents:
            findings.append("Historical incident matches:")
            for incident in incidents:
                if isinstance(incident, dict):
                    findings.append(f"Incident {incident.get('incident_id')} (similarity {incident.get('similarity')})")

        root_causes = state.get("root_causes") or []
        if root_causes:
            findings.append("Root cause hypotheses:")
            for cause in root_causes:
                if isinstance(cause, dict):
                    findings.append(f"{cause.get('cause')} ({cause.get('confidence', 0)}% confidence)")

        self._add_bullet_slide(presentation, "Key Findings", findings)

        recommendations = state.get("recommendations") or []
        if recommendations:
            self._add_bullet_slide(presentation, "Recommendations", [str(item) for item in recommendations])
        else:
            self._add_bullet_slide(presentation, "Recommendations", ["No recommendations were generated."])

        buffer = io.BytesIO()
        presentation.save(buffer)
        buffer.seek(0)
        return buffer.read()

    def _add_bullet_slide(self, presentation: Presentation, title: str, lines: list[str]) -> None:
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        body.text = lines[0] if lines else ""
        for line in lines[1:]:
            paragraph = body.add_paragraph()
            paragraph.level = 1
            paragraph.text = line
            paragraph.font.size = Pt(14)
